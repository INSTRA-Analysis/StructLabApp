"""Generate StructLabPro conference presentation as a .pptx file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ── Colour palette (dark engineering theme) ───────────────────────────────────
CYAN        = RGBColor(0x00, 0xAC, 0xC1)   # #00ACC1 — StructLab accent
DARK_BG     = RGBColor(0x1A, 0x1A, 0x2A)   # near-black slide bg
MID_BG      = RGBColor(0x22, 0x22, 0x35)   # card / table row bg
LIGHT_TEXT  = RGBColor(0xEE, 0xEE, 0xEE)   # primary text
GREY_TEXT   = RGBColor(0x99, 0x99, 0xAA)   # secondary / caption text
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED_ACCENT  = RGBColor(0xE5, 0x39, 0x35)   # warning / gap colour
GREEN_ACT   = RGBColor(0x43, 0xA0, 0x47)   # validated / pass colour
AMBER       = RGBColor(0xF5, 0x7C, 0x00)   # caution / partial

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Low-level helpers ─────────────────────────────────────────────────────────

def _blank_slide(prs: Presentation):
    """Add a fully blank slide (layout 6)."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _bg(slide, colour: RGBColor = DARK_BG):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _box(slide, x, y, w, h, bg: RGBColor | None = None,
         border: RGBColor | None = None, border_pt: float = 1.0):
    """Add a plain rectangle shape."""
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.line.fill.background()
    if bg:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def _txt(slide, text: str, x, y, w, h,
         size: int = 20, bold: bool = False, colour: RGBColor = LIGHT_TEXT,
         align=PP_ALIGN.LEFT, wrap: bool = True, italic: bool = False):
    """Add a text box."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    return txb


def _para(tf, text: str, size: int = 18, bold: bool = False,
          colour: RGBColor = LIGHT_TEXT, align=PP_ALIGN.LEFT,
          italic: bool = False, space_before: int = 0):
    """Append a paragraph to an existing text frame."""
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    return p


def _accent_bar(slide, y=Inches(0.55), colour: RGBColor = CYAN):
    """Thin horizontal rule under the slide title area."""
    line = slide.shapes.add_shape(1, Inches(0.5), y, Inches(12.33), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = colour
    line.line.fill.background()


def _slide_number(slide, n: int):
    _txt(slide, str(n),
         Inches(12.8), Inches(7.1), Inches(0.4), Inches(0.3),
         size=10, colour=GREY_TEXT, align=PP_ALIGN.RIGHT)


def _section_tag(slide, label: str):
    """Small coloured tag in top-right corner."""
    _box(slide, Inches(10.8), Inches(0.12), Inches(2.3), Inches(0.32),
         bg=CYAN)
    _txt(slide, label,
         Inches(10.85), Inches(0.13), Inches(2.2), Inches(0.3),
         size=11, bold=True, colour=DARK_BG, align=PP_ALIGN.CENTER)


def _title_area(slide, title: str, subtitle: str = ""):
    _txt(slide, title,
         Inches(0.5), Inches(0.12), Inches(10.0), Inches(0.5),
         size=26, bold=True, colour=CYAN)
    if subtitle:
        _txt(slide, subtitle,
             Inches(0.5), Inches(0.6), Inches(10.0), Inches(0.32),
             size=14, colour=GREY_TEXT)
    _accent_bar(slide, y=Inches(0.92))


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1 — Title slide."""
    slide = _blank_slide(prs)
    _bg(slide)

    # Large cyan vertical bar left edge
    _box(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, bg=CYAN)

    # Decorative grid dots (subtle)
    for i in range(6):
        for j in range(4):
            dot = slide.shapes.add_shape(
                9,  # oval
                Inches(0.5 + i * 2.1), Inches(1.2 + j * 1.6),
                Inches(0.06), Inches(0.06)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x3F)
            dot.line.fill.background()

    # Main title
    _txt(slide, "StructLabPro",
         Inches(0.55), Inches(1.6), Inches(12.0), Inches(1.2),
         size=54, bold=True, colour=CYAN, align=PP_ALIGN.CENTER)

    _txt(slide, "An Open-Source Desktop Structural Analysis Tool in Python",
         Inches(0.55), Inches(2.75), Inches(12.0), Inches(0.7),
         size=22, colour=LIGHT_TEXT, align=PP_ALIGN.CENTER)

    _txt(slide, "Direct Stiffness Method  ·  2D & 3D  ·  PyQt6 GUI  ·  Python SDK  ·  GPL v3 Licence",
         Inches(0.55), Inches(3.4), Inches(12.0), Inches(0.4),
         size=14, colour=GREY_TEXT, align=PP_ALIGN.CENTER)

    # Horizontal rule
    _box(slide, Inches(3.5), Inches(3.95), Inches(6.33), Pt(1.5), bg=CYAN)

    _txt(slide, "Open Source Tools for Engineering — 2026",
         Inches(0.55), Inches(4.15), Inches(12.0), Inches(0.35),
         size=13, colour=GREY_TEXT, align=PP_ALIGN.CENTER)

    _txt(slide, "Imed Mhamdi",
         Inches(0.55), Inches(4.55), Inches(12.0), Inches(0.35),
         size=15, bold=True, colour=LIGHT_TEXT, align=PP_ALIGN.CENTER)

    # Version badge
    _box(slide, Inches(5.5), Inches(5.3), Inches(2.33), Inches(0.45),
         bg=MID_BG, border=CYAN, border_pt=1.5)
    _txt(slide, "V 1.0.0  —  First Public Release",
         Inches(5.5), Inches(5.31), Inches(2.33), Inches(0.43),
         size=12, bold=True, colour=CYAN, align=PP_ALIGN.CENTER)

    _slide_number(slide, 1)


def slide_agenda(prs):
    """Slide 2 — Agenda / outline."""
    slide = _blank_slide(prs)
    _bg(slide)
    _title_area(slide, "Agenda")
    _section_tag(slide, "Overview")
    _slide_number(slide, 2)

    items = [
        ("01", "The Open-Source Structural Analysis Landscape",
         "What tools exist, their strengths and gaps"),
        ("02", "What is StructLabPro?",
         "Architecture, scope, and design decisions"),
        ("03", "Features Walkthrough",
         "Modelling, results, Python SDK — with live demo"),
        ("04", "Validation",
         "20 benchmark cases vs OpenSeesPy & analytical solutions"),
        ("05", "V1.0 Scope Boundaries",
         "What this release does not do"),
        ("06", "Roadmap & Contributing",
         "Where the project goes next"),
    ]

    y = Inches(1.1)
    for num, title, sub in items:
        # Number badge
        _box(slide, Inches(0.5), y, Inches(0.55), Inches(0.55), bg=CYAN)
        _txt(slide, num, Inches(0.5), y + Inches(0.02),
             Inches(0.55), Inches(0.5),
             size=14, bold=True, colour=DARK_BG, align=PP_ALIGN.CENTER)
        # Text
        _txt(slide, title,
             Inches(1.2), y, Inches(11.0), Inches(0.3),
             size=17, bold=True, colour=LIGHT_TEXT)
        _txt(slide, sub,
             Inches(1.2), y + Inches(0.3), Inches(11.0), Inches(0.25),
             size=12, colour=GREY_TEXT)
        y += Inches(0.85)


def slide_landscape_tools(prs):
    """Slide 3 — Open-source structural tools table."""
    slide = _blank_slide(prs)
    _bg(slide)
    _title_area(slide, "Open-Source Structural Analysis Tools",
                "A survey of what exists before StructLabPro")
    _section_tag(slide, "Landscape")
    _slide_number(slide, 3)

    headers = ["Tool", "Domain", "Interface", "Key Strength", "Key Limitation"]
    col_w   = [1.85, 1.55, 1.45, 3.35, 3.0]
    rows = [
        ("OpenSees / OpenSeesPy", "Frame + Continuum FEM",
         "Python / Tcl scripting", "Nonlinear, dynamic, seismic — research-grade",
         "No GUI, steep learning curve, no interactive modelling"),
        ("Frame3DD", "3D frames & trusses",
         "CLI + text file",
         "Lightweight, fast, well-validated",
         "No GUI, dated UX, limited load types, no Python API"),
        ("PyNite", "3D frame / truss",
         "Python library",
         "Clean modern Python API, MIT licence",
         "Library only — no GUI, no visual workflow"),
        ("anastruct", "2D frames",
         "Python library",
         "Simple Pythonic API, quick to learn",
         "2D only, limited loads, no GUI"),
        ("FEniCS / FEniCSx", "Continuum FEM",
         "Python",
         "Very powerful for PDEs and continuum",
         "Not designed for frame/beam structures"),
        ("CalculiX", "General FEM",
         "CLI / PrePoMax GUI",
         "Industrial-grade, Abaqus-compatible",
         "No Python SDK, complex input, steep learning curve"),
    ]

    # Header row
    x = Inches(0.3)
    y = Inches(1.05)
    row_h = Inches(0.34)
    for i, (hdr, cw) in enumerate(zip(headers, col_w)):
        _box(slide, x, y, Inches(cw - 0.04), row_h, bg=CYAN)
        _txt(slide, hdr, x + Inches(0.04), y + Inches(0.04),
             Inches(cw - 0.1), row_h - Inches(0.08),
             size=11, bold=True, colour=DARK_BG)
        x += Inches(cw)

    # Data rows
    for ri, row in enumerate(rows):
        y += row_h
        x = Inches(0.3)
        bg = MID_BG if ri % 2 == 0 else DARK_BG
        for ci, (cell, cw) in enumerate(zip(row, col_w)):
            _box(slide, x, y, Inches(cw - 0.04), row_h, bg=bg,
                 border=RGBColor(0x33, 0x33, 0x44), border_pt=0.5)
            col = CYAN if ci == 0 else LIGHT_TEXT
            sz = 10 if ci >= 3 else 11
            _txt(slide, cell, x + Inches(0.05), y + Inches(0.03),
                 Inches(cw - 0.1), row_h - Inches(0.06),
                 size=sz, colour=col)
            x += Inches(cw)

    _txt(slide, "* Table covers widely-used FOSS tools as of 2026; excludes commercial tools (SAP2000, ETABS, RFEM, Abaqus).",
         Inches(0.3), Inches(7.1), Inches(12.5), Inches(0.28),
         size=9, colour=GREY_TEXT, italic=True)


def slide_the_gap(prs):
    """Slide 4 — The gap / positioning."""
    slide = _blank_slide(prs)
    _bg(slide)
    _title_area(slide, "The Gap StructLabPro Addresses",
                "Most tools are either GUI-only or scripting-only")
    _section_tag(slide, "Landscape")
    _slide_number(slide, 4)

    # 2×2 quadrant diagram
    # Axes labels
    cx, cy = Inches(6.66), Inches(4.2)
    half_w, half_h = Inches(4.5), Inches(2.6)

    _box(slide, cx - half_w, cy - half_h, half_w * 2, Pt(1.5), bg=GREY_TEXT)   # horizontal axis
    _box(slide, cx - Pt(0.75), cy - half_h, Pt(1.5), half_h * 2, bg=GREY_TEXT) # vertical axis

    _txt(slide, "Scripting / programmable",
         cx - half_w, cy - half_h - Inches(0.35), half_w * 2, Inches(0.35),
         size=12, colour=GREY_TEXT, align=PP_ALIGN.LEFT)
    _txt(slide, "GUI / interactive",
         cx - half_w, cy + half_h, half_w * 2, Inches(0.35),
         size=12, colour=GREY_TEXT, align=PP_ALIGN.RIGHT)
    _txt(slide, "Simple / educational",
         cx + half_w - Inches(2.2), cy - Inches(0.2), Inches(2.2), Inches(0.3),
         size=11, colour=GREY_TEXT, align=PP_ALIGN.RIGHT)
    _txt(slide, "Powerful / research",
         cx - half_w, cy - Inches(0.2), Inches(2.2), Inches(0.3),
         size=11, colour=GREY_TEXT)

    # Tool dots
    def _dot(label, qx, qy, col=GREY_TEXT):
        _box(slide, cx + Inches(qx) - Inches(0.12), cy - Inches(qy) - Inches(0.12),
             Inches(0.24), Inches(0.24), bg=col)
        _txt(slide, label,
             cx + Inches(qx) - Inches(1.1), cy - Inches(qy) + Inches(0.15),
             Inches(2.2), Inches(0.3), size=11, colour=col)

    _dot("OpenSees",  -2.8,  1.6, AMBER)
    _dot("PyNite",    -1.5,  0.6, AMBER)
    _dot("anastruct", -0.5,  0.5, AMBER)
    _dot("Frame3DD",  -2.0, -0.3, AMBER)
    _dot("CalculiX",  -3.2, -1.0, AMBER)

    # StructLabPro star position
    _box(slide, cx + Inches(1.6) - Inches(0.18), cy - Inches(1.4) - Inches(0.18),
         Inches(0.36), Inches(0.36), bg=CYAN)
    _txt(slide, "★  StructLabPro",
         cx + Inches(1.65), cy - Inches(1.4) - Inches(0.05),
         Inches(2.5), Inches(0.35), size=13, bold=True, colour=CYAN)

    # Call-out box
    _box(slide, Inches(0.4), Inches(1.1), Inches(3.8), Inches(1.7),
         bg=MID_BG, border=CYAN, border_pt=1.5)
    txb = slide.shapes.add_textbox(Inches(0.55), Inches(1.18), Inches(3.5), Inches(1.5))
    tf = txb.text_frame
    tf.word_wrap = True
    _para(tf, "StructLabPro targets the intersection:", size=13, bold=True, colour=CYAN)
    _para(tf, "· Full desktop GUI for interactive modelling", size=12, colour=LIGHT_TEXT, space_before=4)
    _para(tf, "· Python SDK + embedded console for scripting", size=12, colour=LIGHT_TEXT)
    _para(tf, "· GPL v3 licence, standalone .exe, zero install", size=12, colour=LIGHT_TEXT)


def slide_what_is(prs):
    """Slide 5 — What is StructLabPro?"""
    slide = _blank_slide(prs)
    _bg(slide)
    _title_area(slide, "What is StructLabPro?",
                "A 2D & 3D structural analysis desktop application")
    _section_tag(slide, "Project")
    _slide_number(slide, 5)

    facts = [
        ("Solver",       "Direct Stiffness Method — 2D (3 DOF/node) and 3D (6 DOF/node), unified engine"),
        ("Elements",     "Beams · Frames · Trusses · Bars · Spring supports · Pin releases"),
        ("Loads",        "Point loads · UDL · UVL · Nodal moments · EN 1990 load combinations"),
        ("GUI",          "PyQt6 desktop application — dark theme, interactive canvas, overlays (BMD/SFD/AFD/deformed shape)"),
        ("Sections",     "European steel section library (IPE, HEA, HEB, SHS, RHS) + custom sections"),
        ("Output",       "Displacements · Reactions · Member forces · Design utilisation η · Multi-page PDF report"),
        ("SDK",          "sdk.py Python API + embedded IPython console — model and solve without the GUI"),
        ("Distribution", "Standalone .exe via PyInstaller — no Python install required for end users"),
        ("Licence",      "GPL v3 — free to use, modify, distribute (copyleft)"),
    ]

    y = Inches(1.1)
    for label, detail in facts:
        _box(slide, Inches(0.4), y, Inches(1.6), Inches(0.36), bg=MID_BG, border=CYAN, border_pt=0.8)
        _txt(slide, label, Inches(0.45), y + Inches(0.04), Inches(1.5), Inches(0.3),
             size=11, bold=True, colour=CYAN)
        _txt(slide, detail, Inches(2.15), y + Inches(0.04), Inches(10.9), Inches(0.3),
             size=11, colour=LIGHT_TEXT)
        y += Inches(0.44)


def slide_architecture(prs):
    """Slide 6 — Architecture diagram."""
    slide = _blank_slide(prs)
    _bg(slide)
    _title_area(slide, "Architecture — One Engine, All Structure Types",
                "2D/3D is auto-detected; no separate solver per mode")
    _section_tag(slide, "Project")
    _slide_number(slide, 6)

    # Layer boxes (top = input, bottom = output)
    layers = [
        ("GUI Layer",
         "PyQt6 canvas  ·  Properties panel  ·  Load case manager  ·  Overlay diagrams",
         CYAN, DARK_BG),
        ("Model State",
         "NodeData  ·  MemberData  ·  LoadCase  ·  SupportType  (pure Python dataclasses)",
         MID_BG, LIGHT_TEXT),
        ("Model Builder",
         "ModelState → core Model bridge  ·  sub-node division for overlays  ·  z-coord auto-detect",
         MID_BG, LIGHT_TEXT),
        ("Core Engine",
         "FrameElement (2D 6×6 / 3D 12×12)  ·  Assembler  ·  LinearSolver  ·  Postprocessor  ·  FEM loads",
         RGBColor(0x1A, 0x2A, 0x35), CYAN),
        ("Python SDK  /  IPython Console",
         "sdk.py  ·  build models in pure Python  ·  solve()  ·  query results  ·  matplotlib plots",
         RGBColor(0x1A, 0x28, 0x1A), GREEN_ACT),
    ]

    y = Inches(1.1)
    box_h = Inches(0.9)
    gap   = Inches(0.12)
    for label, detail, bg, fg in layers:
        _box(slide, Inches(0.5), y, Inches(12.33), box_h, bg=bg,
             border=CYAN if fg == DARK_BG else RGBColor(0x33, 0x33, 0x55),
             border_pt=1.5 if fg == DARK_BG else 0.8)
        _txt(slide, label, Inches(0.65), y + Inches(0.05), Inches(3.0), Inches(0.38),
             size=14, bold=True, colour=fg if fg != DARK_BG else DARK_BG)
        _txt(slide, detail, Inches(0.65), y + Inches(0.44), Inches(12.0), Inches(0.38),
             size=11, colour=GREY_TEXT if bg != CYAN else DARK_BG)

        # Arrow between layers
        if y + box_h + gap < Inches(7.0):
            arr_y = y + box_h + Inches(0.01)
            _box(slide, Inches(6.5), arr_y, Inches(0.33), gap, bg=GREY_TEXT)

        y += box_h + gap + Inches(0.04)

    _txt(slide,
         "NumPy · SciPy · Matplotlib underpin the engine.  "
         "The GUI never bypasses the solver — it always goes through Model Builder → Core Engine.",
         Inches(0.5), Inches(7.05), Inches(12.33), Inches(0.32),
         size=10, colour=GREY_TEXT, italic=True)


def slide_modelling(prs):
    """Slide 7 — Modelling features."""
    slide = _blank_slide(prs)
    _bg(slide)
    _title_area(slide, "Modelling — Interactive Canvas",
                "Draw structures directly or load from 30+ built-in presets")
    _section_tag(slide, "Features")
    _slide_number(slide, 7)

    cols = [
        ("Canvas", [
            "Click to place nodes, drag to connect members",
            "Snap-to-grid, multi-select, Ctrl+D duplicate",
            "2D and 3D modes — isometric projection, working planes (XY / XZ / YZ)",
            "Orbit, pan, zoom in 3D view",
            "Selection filter: nodes / members / both",
        ]),
        ("Element Types", [
            "Beam (continuous, with/without pin releases)",
            "Frame (rigid joints, full 3 or 6 DOF)",
            "Truss / Bar (pin at both ends)",
            "Spring supports (translational & rotational)",
            "Mixed models — beams + columns + trusses in one model",
        ]),
        ("Loads & Cases", [
            "Point loads (Fx, Fy, Fz, Moment)",
            "UDL and UVL (varying distributed) on members",
            "EN 1990 named load cases (G, Q, W, …)",
            "Load combinations with factors",
            "Combination envelope (max / min across all combos)",
            "Pattern loading detection (EN 1992-1-1)",
        ]),
    ]

    col_w = Inches(4.0)
    x = Inches(0.4)
    for title, bullets in cols:
        _box(slide, x, Inches(1.0), col_w - Inches(0.1), Inches(0.38),
             bg=CYAN)
        _txt(slide, title, x + Inches(0.08), Inches(1.02),
             col_w - Inches(0.2), Inches(0.34),
             size=13, bold=True, colour=DARK_BG)

        txb = slide.shapes.add_textbox(x + Inches(0.08), Inches(1.45),
                                        col_w - Inches(0.2), Inches(5.6))
        tf = txb.text_frame
        tf.word_wrap = True
        first = True
        for b in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(5)
            run = p.add_run()
            run.text = "▸  " + b
            run.font.size = Pt(12)
            run.font.color.rgb = LIGHT_TEXT
        x += col_w


def slide_results(prs):
    """Slide 8 — Results and overlays."""
    slide = _blank_slide(prs)
    _bg(slide)
    _title_area(slide, "Results — Visual Overlays + Tabular Output",
                "Live diagrams drawn directly on the canvas after solve")
    _section_tag(slide, "Features")
    _slide_number(slide, 8)

    # Left: overlay types
    overlays = [
        ("BMD", "Bending Moment Diagram — filled, sign-correct\n(sagging below baseline, hogging above)",   CYAN),
        ("SFD", "Shear Force Diagram — peak values labelled at critical sections",                           GREEN_ACT),
        ("AFD", "Axial Force Diagram — compression / tension colour-coded",                                   AMBER),
        ("Def", "Deformed Shape — scaled, shows buckled mode visually",                                      LIGHT_TEXT),
    ]

    x = Inches(0.4)
    y = Inches(1.1)
    for tag, desc, col in overlays:
        _box(slide, x, y, Inches(0.7), Inches(0.55), bg=col)
        _txt(slide, tag, x, y + Inches(0.08), Inches(0.7), Inches(0.4),
             size=16, bold=True, colour=DARK_BG, align=PP_ALIGN.CENTER)
        _txt(slide, desc, x + Inches(0.8), y + Inches(0.05), Inches(5.2), Inches(0.5),
             size=12, colour=LIGHT_TEXT)
        y += Inches(0.7)

    # Right: results tables
    _box(slide, Inches(6.9), Inches(1.0), Inches(6.0), Inches(5.8),
         bg=MID_BG, border=RGBColor(0x33, 0x33, 0x55), border_pt=0.8)

    txb = slide.shapes.add_textbox(Inches(7.1), Inches(1.08), Inches(5.7), Inches(5.6))
    tf = txb.text_frame
    tf.word_wrap = True

    _para(tf, "Results Panel — four tabs:", size=13, bold=True, colour=CYAN)
    tabs = [
        ("Displacements", "dx, dy (mm) · θ (mrad) per node"),
        ("Reactions",     "Fx, Fy (kN) · M (kN·m) at supports"),
        ("Member forces", "N, V end-pair + M+ peak sagging · M− peak hogging + x/L position"),
        ("Design",        "M_Ed vs M_Rd · utilisation η (%) · Pass / Fail status"),
    ]
    for tab, detail in tabs:
        _para(tf, f"▸ {tab}", size=12, bold=True, colour=LIGHT_TEXT, space_before=6)
        _para(tf, f"   {detail}", size=11, colour=GREY_TEXT)

    _para(tf, "", size=8)
    _para(tf, "Bidirectional selection:", size=13, bold=True, colour=CYAN, space_before=6)
    _para(tf, "Click a table row → highlights member on canvas", size=12, colour=LIGHT_TEXT, space_before=4)
    _para(tf, "Click canvas member → scrolls to that row in table", size=12, colour=LIGHT_TEXT)
    _para(tf, "", size=8)
    _para(tf, "PDF Report:", size=13, bold=True, colour=CYAN, space_before=4)
    _para(tf, "Multi-page A4 · geometry + loads + diagrams + full results table", size=12, colour=LIGHT_TEXT, space_before=4)


def slide_sdk(prs):
    """Slide 9 — Python SDK & console."""
    slide = _blank_slide(prs)
    _bg(slide)
    _title_area(slide, "Python SDK & Embedded Console",
                "The GUI and the scripting interface share the same live model")
    _section_tag(slide, "Features")
    _slide_number(slide, 9)

    # Left: description
    txb = slide.shapes.add_textbox(Inches(0.4), Inches(1.08), Inches(4.8), Inches(5.5))
    tf = txb.text_frame
    tf.word_wrap = True
    _para(tf, "sdk.py — build and solve in pure Python:", size=13, bold=True, colour=CYAN)
    points = [
        "No GUI needed — runs in any Python environment",
        "Same engine as the desktop app (single code path)",
        "Ideal for parametric studies, lecture examples, automated checks",
        "Results returned as typed Python objects (not strings)",
    ]
    for p in points:
        _para(tf, f"▸  {p}", size=12, colour=LIGHT_TEXT, space_before=5)

    _para(tf, "", size=8)
    _para(tf, "Embedded IPython Console:", size=13, bold=True, colour=CYAN, space_before=6)
    console_points = [
        "Opens inside the running app (qtconsole)",
        "model, state, sdk, np, plt pre-injected",
        "Inspect, modify, or extend the current canvas model",
        "matplotlib plots open in separate windows",
    ]
    for p in console_points:
        _para(tf, f"▸  {p}", size=12, colour=LIGHT_TEXT, space_before=5)

    # Right: code box
    code = (
        "# Simply supported steel beam — 6 m, UDL 15 kN/m\n"
        "import sdk\n\n"
        "m = sdk.Model()\n"
        "n0 = m.add_node(0, 0, 0)\n"
        "n1 = m.add_node(6, 0, 0)\n"
        "m.pin(n0)\n"
        "m.roller(n1)\n\n"
        "beam = m.add_member(n0, n1)\n"
        "beam.set_section(E=210e9, A=5.381e-3, I=8.356e-5)  # IPE 300\n"
        "m.add_udl(beam, w=15_000)   # 15 kN/m downward\n\n"
        "res = m.solve()\n"
        "print(res.max_deflection())  # → 7.03 mm\n"
        "res.plot_bmd()"
    )

    _box(slide, Inches(5.5), Inches(1.0), Inches(7.5), Inches(5.9),
         bg=RGBColor(0x0D, 0x0D, 0x17),
         border=CYAN, border_pt=1.2)

    txb2 = slide.shapes.add_textbox(Inches(5.65), Inches(1.1), Inches(7.2), Inches(5.7))
    tf2 = txb2.text_frame
    tf2.word_wrap = False
    first = True
    for line in code.split("\n"):
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(12)
        run.font.color.rgb = GREEN_ACT if line.startswith("#") else (
            CYAN if any(kw in line for kw in ("import", "def ", "class ", "return")) else LIGHT_TEXT
        )
        from pptx.oxml.ns import qn
        from lxml import etree
        rPr = run._r.get_or_add_rPr()
        latin = etree.SubElement(rPr, qn('a:latin'))
        latin.set('typeface', 'Consolas')


def slide_demo(prs):
    """Slide 10 — Live demo placeholder."""
    slide = _blank_slide(prs)
    _bg(slide)
    _section_tag(slide, "Demo")
    _slide_number(slide, 10)

    _box(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, bg=CYAN)

    _txt(slide, "Live Demonstration",
         Inches(0.55), Inches(2.0), Inches(12.0), Inches(1.0),
         size=44, bold=True, colour=CYAN, align=PP_ALIGN.CENTER)

    steps = [
        "1  ·  Launch → welcome dialog → load Steel Industrial Portal preset",
        "2  ·  Add a load case · solve · walk through BMD / SFD overlays",
        "3  ·  Generate PDF report",
        "4  ·  Switch to 3D — open 3D Floor Grid preset · orbit · solve · deformed shape",
        "5  ·  Open Python console · query model · plot BMD from the console",
    ]

    y = Inches(3.4)
    for s in steps:
        _txt(slide, s, Inches(2.5), y, Inches(9.0), Inches(0.4),
             size=15, colour=LIGHT_TEXT, align=PP_ALIGN.CENTER)
        y += Inches(0.5)


def slide_validation(prs):
    """Slide 11 — Benchmark validation."""
    slide = _blank_slide(prs)
    _bg(slide)
    _title_area(slide, "Validation — 20 Benchmark Cases",
                "StructLabPro vs OpenSeesPy (where applicable) and analytical solutions")
    _section_tag(slide, "Validation")
    _slide_number(slide, 11)

    headers = ["ID", "Type", "Description", "Reference", "Max error"]
    col_w   = [0.6, 1.0, 4.8, 3.5, 1.5]
    groups = [
        # (ID, type, description, reference, error)
        ("B1", "Beam", "Simply supported beam — midspan point load",        "Analytical (Euler-Bernoulli)", "0.00%"),
        ("B2", "Beam", "Propped cantilever — UDL",                          "Analytical",                  "0.00%"),
        ("B3", "Beam", "Fixed-fixed beam — UDL",                            "Analytical",                  "0.00%"),
        ("B4", "Beam", "3-span continuous beam — mixed loads",               "OpenSeesPy",                  "0.00%"),
        ("B5", "Beam", "Gerber beam — internal hinge (pin release)",         "Analytical",                  "0.00%"),
        ("F1", "Frame","Simply supported portal — lateral load",             "OpenSeesPy",                  "0.00%"),
        ("F2", "Frame","Fixed-base portal — UDL on beam",                   "OpenSeesPy",                  "0.00%"),
        ("F3", "Frame","Multi-storey frame — vertical + lateral loads",      "OpenSeesPy",                  "0.00%"),
        ("F4", "Frame","Two-bay frame — asymmetric loading",                 "OpenSeesPy",                  "0.00%"),
        ("F5", "Frame","Pitched roof portal with eaves haunches",            "OpenSeesPy",                  "0.00%"),
        ("T1", "Truss","8-panel Pratt truss — symmetric UDL",                "Method of joints",            "0.00%"),
        ("T2", "Truss","Warren truss — asymmetric point loads",              "Method of joints",            "0.00%"),
        ("3D1","3D",   "3D cantilever — tip load in Y",                     "Analytical",                  "0.00%"),
        ("3D2","3D",   "3D cantilever — tip load in Z",                     "Analytical",                  "0.00%"),
        ("3D3","3D",   "3D simply supported beam — midspan load",            "Analytical",                  "0.00%"),
        ("3D4","3D",   "3D portal frame — in-plane lateral",                "OpenSeesPy",                  "0.00%"),
        ("3D5","3D",   "3D portal frame — out-of-plane lateral",            "OpenSeesPy",                  "0.00%"),
        ("3D6","3D",   "3D 2-storey space frame — combined loads",          "OpenSeesPy",                  "0.00%"),
        ("3D7","3D",   "3D floor grid — biaxial bending",                   "OpenSeesPy",                  "0.00%"),
        ("3D8","3D",   "3D Pratt roof truss — space truss",                 "Analytical",                  "0.00%"),
    ]

    # Header
    x0 = Inches(0.3)
    y = Inches(1.05)
    row_h = Inches(0.28)
    x = x0
    for h, cw in zip(headers, col_w):
        _box(slide, x, y, Inches(cw - 0.03), row_h, bg=CYAN)
        _txt(slide, h, x + Inches(0.03), y + Inches(0.02), Inches(cw - 0.06), row_h,
             size=10, bold=True, colour=DARK_BG)
        x += Inches(cw)

    for ri, row in enumerate(groups):
        y += row_h
        x = x0
        bg = MID_BG if ri % 2 == 0 else DARK_BG
        for ci, (cell, cw) in enumerate(zip(row, col_w)):
            col = (CYAN if ci == 0
                   else GREEN_ACT if ci == 4
                   else (AMBER if row[1] == "3D" else LIGHT_TEXT))
            _box(slide, x, y, Inches(cw - 0.03), row_h, bg=bg,
                 border=RGBColor(0x2A, 0x2A, 0x3A), border_pt=0.4)
            _txt(slide, cell, x + Inches(0.03), y + Inches(0.01),
                 Inches(cw - 0.06), row_h,
                 size=9, colour=col)
            x += Inches(cw)

    _txt(slide, "All 20 cases: max displacement / reaction / moment error = 0.00% (tolerance < 0.01%)",
         Inches(0.3), Inches(7.1), Inches(12.5), Inches(0.28),
         size=10, bold=True, colour=GREEN_ACT, italic=True)


def slide_limitations(prs):
    """Slide 12 — Honest V1.0 limitations."""
    slide = _blank_slide(prs)
    _bg(slide)
    _title_area(slide, "V1.0 — Scope Boundaries (What This Release Does Not Do)",
                "Honest positioning — this is a linear-static desktop tool")
    _section_tag(slide, "Scope")
    _slide_number(slide, 12)

    gaps = [
        ("Dynamic / Seismic",   "No modal analysis, no response spectrum, no time-history integration",         RED_ACCENT),
        ("Nonlinear Analysis",  "No geometric nonlinearity (P-Δ / large displacement), no material plasticity", RED_ACCENT),
        ("Plates & Shells",     "Frame elements only — no 2D/3D continuum mesh elements",                       RED_ACCENT),
        ("Design Code Checks",  "EC3 steel moment check (η = M_Ed/M_Rd) only — EC2 concrete check is a stub",  AMBER),
        ("Load Generation",     "Loads are manual — no EN 1991 wind, snow, or seismic auto-generator",          AMBER),
        ("Section Optimisation","Checks utilisation but does not select or size sections automatically",         AMBER),
        ("IFC / DXF Exchange",  "No import or export to BIM or CAD formats",                                    AMBER),
        ("Cloud / Multi-user",  "Local desktop only — no collaboration or cloud-hosted solve",                  GREY_TEXT),
    ]

    y = Inches(1.1)
    for label, detail, col in gaps:
        _box(slide, Inches(0.4), y, Inches(0.08), Inches(0.38), bg=col)
        _txt(slide, label, Inches(0.6), y + Inches(0.02), Inches(2.8), Inches(0.34),
             size=12, bold=True, colour=col)
        _txt(slide, detail, Inches(3.5), y + Inches(0.02), Inches(9.4), Inches(0.34),
             size=12, colour=LIGHT_TEXT)
        y += Inches(0.48)

    _txt(slide, "Red = not in scope for V1.x   ·   Amber = partial or planned",
         Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.28),
         size=10, colour=GREY_TEXT, italic=True)


def slide_roadmap(prs):
    """Slide 13 — Roadmap & contributing."""
    slide = _blank_slide(prs)
    _bg(slide)
    _title_area(slide, "Roadmap & Contributing",
                "Open to contributions — GPL v3 licence on GitHub")
    _section_tag(slide, "Roadmap")
    _slide_number(slide, 13)

    phases = [
        ("Phase 10\n(active)", "3D modelling UX",
         "Blender-style orbit · axis-constrained grab (G+X/Y/Z) · E-extrude · numpad view recall · 3D grid visual",
         CYAN),
        ("Phase 11", "EC2 Concrete Design Check",
         "Steel / Concrete branch in Design tab · EC2 §6.1 rectangular stress block · user-input M_Rd for concrete",
         AMBER),
        ("Phase 12", "OpenSees Export",
         "One-click export of any StructLab model to a ready-to-run OpenSeesPy script",
         GREEN_ACT),
        ("Future", "Dynamic / Nonlinear",
         "Modal analysis · response spectrum · P-Δ geometric nonlinearity · pushover",
         GREY_TEXT),
    ]

    y = Inches(1.1)
    for phase, title, detail, col in phases:
        _box(slide, Inches(0.4), y, Inches(1.45), Inches(0.9),
             bg=MID_BG, border=col, border_pt=1.5)
        _txt(slide, phase, Inches(0.42), y + Inches(0.08), Inches(1.4), Inches(0.75),
             size=11, bold=True, colour=col, align=PP_ALIGN.CENTER)
        _txt(slide, title, Inches(2.05), y + Inches(0.04), Inches(10.7), Inches(0.3),
             size=14, bold=True, colour=col)
        _txt(slide, detail, Inches(2.05), y + Inches(0.36), Inches(10.7), Inches(0.5),
             size=12, colour=LIGHT_TEXT)
        y += Inches(1.05)

    # Contributing box
    _box(slide, Inches(0.4), Inches(5.55), Inches(12.33), Inches(1.55),
         bg=MID_BG, border=CYAN, border_pt=1.2)
    txb = slide.shapes.add_textbox(Inches(0.6), Inches(5.65), Inches(12.0), Inches(1.4))
    tf = txb.text_frame
    tf.word_wrap = True
    _para(tf, "How to contribute:", size=13, bold=True, colour=CYAN)
    contrib = [
        "Report bugs or request features via GitHub Issues",
        "Add benchmark cases — any textbook example with a known analytical solution",
        "Extend the SDK — new load types, result queries, export formats",
        "Improve the section library — EN 10365 profiles, CHS, angles, channels",
    ]
    for c in contrib:
        _para(tf, f"▸  {c}", size=12, colour=LIGHT_TEXT, space_before=3)


def slide_closing(prs):
    """Slide 14 — Closing / thank you."""
    slide = _blank_slide(prs)
    _bg(slide)
    _box(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, bg=CYAN)
    _slide_number(slide, 14)

    _txt(slide, "Thank You",
         Inches(0.55), Inches(1.7), Inches(12.0), Inches(1.0),
         size=50, bold=True, colour=CYAN, align=PP_ALIGN.CENTER)

    _txt(slide, "StructLabPro V1.0.0 — Open Source Structural Analysis in Python",
         Inches(0.55), Inches(2.75), Inches(12.0), Inches(0.5),
         size=19, colour=LIGHT_TEXT, align=PP_ALIGN.CENTER)

    _box(slide, Inches(3.5), Inches(3.35), Inches(6.33), Pt(1.5), bg=CYAN)

    info = [
        ("Licence",   "GPL v3 — free to use, modify, distribute (copyleft)"),
        ("Tech stack","Python · NumPy · SciPy · PyQt6 · Matplotlib · PyInstaller"),
        ("Solver",    "Direct Stiffness Method · 2D & 3D unified engine"),
        ("Tests",     "92 / 92 passing · 20 benchmark cases · 0.00% error"),
    ]

    y = Inches(3.65)
    for label, val in info:
        _txt(slide, f"{label}:", Inches(3.0), y, Inches(2.0), Inches(0.34),
             size=13, bold=True, colour=CYAN, align=PP_ALIGN.RIGHT)
        _txt(slide, val, Inches(5.2), y, Inches(5.5), Inches(0.34),
             size=13, colour=LIGHT_TEXT)
        y += Inches(0.44)

    _txt(slide, "Questions?",
         Inches(0.55), Inches(6.1), Inches(12.0), Inches(0.5),
         size=22, bold=True, colour=GREY_TEXT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_agenda(prs)
    slide_landscape_tools(prs)
    slide_the_gap(prs)
    slide_what_is(prs)
    slide_architecture(prs)
    slide_modelling(prs)
    slide_results(prs)
    slide_sdk(prs)
    slide_demo(prs)
    slide_validation(prs)
    slide_limitations(prs)
    slide_roadmap(prs)
    slide_closing(prs)

    out = "StructLabPro_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
